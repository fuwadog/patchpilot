"""File manager – loads and tracks project files in context with support for multiple file types."""
from __future__ import annotations
import glob
import os
from typing import Optional
import mimetypes

from context.manager import ContextManager


class FileManager:
    """
    Responsible for discovering, reading, and tracking project files.
    Supports multiple file formats: code, text, PDF, Word, Excel, etc.
    Delegates context injection to ContextManager.
    """

    def __init__(self, context: ContextManager, max_files: int, default_extensions: list[str]):
        self._ctx = context
        self._max_files = max_files
        self._extensions = default_extensions
        # path -> full raw content (for patch generation etc.)
        self._store: dict[str, str] = {}
        
        # Initialize mimetypes
        mimetypes.init()

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def load(self, path: str) -> tuple[bool, Optional[str]]:
        """Load a single file. Returns (success, error_message)."""
        path = os.path.abspath(path)
        content, err = self._read(path)
        if err:
            return False, err
        self._store[path] = content
        self._ctx.upsert_file(path, content)
        return True, None

    def unload(self, path: str, force: bool = False) -> bool:
        """
        Unload a single file. 
        Returns True if successful, False if skipped (e.g. pinned).
        """
        path = os.path.abspath(path)
        if self._ctx.remove_file(path, force=force):
            self._store.pop(path, None)
            return True
        return False

    def unload_all(self, keep_pinned: bool = True) -> int:
        """Unload all non-pinned files."""
        paths = list(self._store.keys())
        count = 0
        for p in paths:
            if self.unload(p, force=not keep_pinned):
                count += 1
        return count

    def unload_folder(self, folder: str) -> int:
        """Unload all files within a folder."""
        folder = os.path.abspath(folder)
        paths = list(self._store.keys())
        count = 0
        for p in paths:
            if p.startswith(folder):
                if self.unload(p):
                    count += 1
        return count

    def unload_pattern(self, pattern: str) -> int:
        """Unload files matching a glob pattern."""
        import fnmatch
        paths = list(self._store.keys())
        count = 0
        for p in paths:
            if fnmatch.fnmatch(os.path.basename(p), pattern) or fnmatch.fnmatch(p, pattern):
                if self.unload(p):
                    count += 1
        return count

    def get_pin_status(self, path: str) -> bool:
        return self._ctx.is_pinned(path)

    def load_folder(
        self, folder: str, extensions: Optional[list[str]] = None
    ) -> tuple[int, list[str]]:
        """
        Discover and load up to max_files matching files from folder.
        Returns (loaded_count, list_of_errors).
        """
        ext_list = extensions or self._extensions
        discovered = self._discover(folder, ext_list)
        errors: list[str] = []
        loaded = 0
        for path in discovered:
            ok, err = self.load(path)
            if ok:
                loaded += 1
            else:
                errors.append(f"{path}: {err}")
        return loaded, errors

    def get_content(self, path: str) -> Optional[str]:
        return self._store.get(os.path.abspath(path))

    def loaded_paths(self) -> list[str]:
        return list(self._store.keys())

    def is_loaded(self, path: str) -> bool:
        return os.path.abspath(path) in self._store

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _discover(self, folder: str, extensions: list[str]) -> list[str]:
        files: list[str] = []
        for ext in extensions:
            files.extend(glob.glob(os.path.join(folder, "**", ext), recursive=True))
        seen: set[str] = set()
        unique: list[str] = []
        for f in sorted(files):
            af = os.path.abspath(f)
            if af not in seen:
                seen.add(af)
                unique.append(af)
        return unique[: self._max_files]

    def _read(self, path: str) -> tuple[Optional[str], Optional[str]]:
        """
        Read file content based on file type.
        Supports: text files, PDFs, Word docs, Excel files, and more.
        """
        if not os.path.exists(path):
            return None, "File not found."
        
        try:
            # Get file extension
            _, ext = os.path.splitext(path)
            ext = ext.lower()
            
            # Route to appropriate reader based on file type
            if ext in ['.pdf']:
                return self._read_pdf(path)
            elif ext in ['.docx', '.doc']:
                return self._read_word(path)
            elif ext in ['.xlsx', '.xls', '.xlsm']:
                return self._read_excel(path)
            elif ext in ['.csv']:
                return self._read_csv(path)
            elif ext in ['.json']:
                return self._read_json(path)
            elif ext in ['.xml']:
                return self._read_xml(path)
            elif ext in ['.md', '.markdown']:
                return self._read_text(path)
            elif ext in ['.rst']:
                return self._read_text(path)
            elif ext in ['.yaml', '.yml']:
                return self._read_text(path)
            elif ext in ['.toml']:
                return self._read_text(path)
            elif ext in ['.ini', '.cfg', '.conf']:
                return self._read_text(path)
            elif ext in ['.html', '.htm', '.xhtml']:
                return self._read_html(path)
            elif ext in ['.rtf']:
                return self._read_rtf(path)
            elif ext in ['.odt']:
                return self._read_odt(path)
            elif ext in ['.pptx', '.ppt']:
                return self._read_powerpoint(path)
            else:
                # Try to read as text file for code and plain text
                return self._read_text(path)
                
        except Exception as exc:
            return None, f"Read error: {exc}"

    @staticmethod
    def _read_text(path: str) -> tuple[Optional[str], Optional[str]]:
        """Read plain text files."""
        try:
            with open(path, "r", encoding="utf-8") as fh:
                return fh.read(), None
        except UnicodeDecodeError:
            # Try with different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(path, "r", encoding=encoding) as fh:
                        return fh.read(), None
                except:
                    continue
            return None, "Unable to decode file with common encodings"
        except Exception as exc:
            return None, f"Read error: {exc}"

    @staticmethod
    def _read_pdf(path: str) -> tuple[Optional[str], Optional[str]]:
        """Extract text from PDF files."""
        try:
            import PyPDF2
            text_content = []
            with open(path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text_content.append(f"--- Page {page_num + 1} ---\n{page.extract_text()}")
            return "\n\n".join(text_content), None
        except ImportError:
            return None, "PyPDF2 not installed. Install with: pip install PyPDF2"
        except Exception as exc:
            return None, f"PDF read error: {exc}"

    @staticmethod
    def _read_word(path: str) -> tuple[Optional[str], Optional[str]]:
        """Extract text from Word documents (.docx, .doc)."""
        try:
            import docx
            doc = docx.Document(path)
            text_content = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    table_text.append(row_text)
                if table_text:
                    text_content.append("\n[TABLE]\n" + "\n".join(table_text) + "\n[/TABLE]\n")
            
            return "\n\n".join(text_content), None
        except ImportError:
            return None, "python-docx not installed. Install with: pip install python-docx"
        except Exception as exc:
            return None, f"Word document read error: {exc}"

    @staticmethod
    def _read_excel(path: str) -> tuple[Optional[str], Optional[str]]:
        """Extract data from Excel files (.xlsx, .xls)."""
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(path, data_only=True)
            text_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"=== Sheet: {sheet_name} ===")
                
                for row in sheet.iter_rows(values_only=True):
                    # Filter out completely empty rows
                    if any(cell is not None for cell in row):
                        row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                        text_content.append(row_text)
                text_content.append("")  # Empty line between sheets
            
            return "\n".join(text_content), None
        except ImportError:
            return None, "openpyxl not installed. Install with: pip install openpyxl"
        except Exception as exc:
            return None, f"Excel read error: {exc}"

    @staticmethod
    def _read_csv(path: str) -> tuple[Optional[str], Optional[str]]:
        """Read CSV files."""
        try:
            import csv
            text_content = []
            with open(path, 'r', encoding='utf-8', newline='') as file:
                csv_reader = csv.reader(file)
                for row in csv_reader:
                    text_content.append(" | ".join(row))
            return "\n".join(text_content), None
        except Exception as exc:
            return None, f"CSV read error: {exc}"

    @staticmethod
    def _read_json(path: str) -> tuple[Optional[str], Optional[str]]:
        """Read and pretty-print JSON files."""
        try:
            import json
            with open(path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                return json.dumps(data, indent=2, ensure_ascii=False), None
        except Exception as exc:
            return None, f"JSON read error: {exc}"

    @staticmethod
    def _read_xml(path: str) -> tuple[Optional[str], Optional[str]]:
        """Read XML files."""
        try:
            import xml.etree.ElementTree as ET
            tree = ET.parse(path)
            root = tree.getroot()
            # Return the XML as formatted string
            return ET.tostring(root, encoding='unicode', method='xml'), None
        except Exception as exc:
            return None, f"XML read error: {exc}"

    @staticmethod
    def _read_html(path: str) -> tuple[Optional[str], Optional[str]]:
        """Extract text from HTML files."""
        try:
            from bs4 import BeautifulSoup
            with open(path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file.read(), 'html.parser')
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                text = soup.get_text()
                # Clean up whitespace
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                text = '\n'.join(chunk for chunk in chunks if chunk)
                return text, None
        except ImportError:
            # Fallback to reading as plain text if BeautifulSoup not available
            return FileManager._read_text(path)
        except Exception as exc:
            return None, f"HTML read error: {exc}"

    @staticmethod
    def _read_rtf(path: str) -> tuple[Optional[str], Optional[str]]:
        """Extract text from RTF files."""
        try:
            from striprtf.striprtf import rtf_to_text
            with open(path, 'r', encoding='utf-8') as file:
                rtf_content = file.read()
                text = rtf_to_text(rtf_content)
                return text, None
        except ImportError:
            return None, "striprtf not installed. Install with: pip install striprtf"
        except Exception as exc:
            return None, f"RTF read error: {exc}"

    @staticmethod
    def _read_odt(path: str) -> tuple[Optional[str], Optional[str]]:
        """Extract text from ODT (OpenDocument Text) files."""
        try:
            from odfdo import Document
            doc = Document(path)
            text_content = []
            for paragraph in doc.body.paragraphs:
                text_content.append(paragraph.text)
            return "\n\n".join(text_content), None
        except ImportError:
            return None, "odfdo not installed. Install with: pip install odfdo"
        except Exception as exc:
            return None, f"ODT read error: {exc}"

    @staticmethod
    def _read_powerpoint(path: str) -> tuple[Optional[str], Optional[str]]:
        """Extract text from PowerPoint files."""
        try:
            from pptx import Presentation
            prs = Presentation(path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content.append(f"=== Slide {slide_num} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
                text_content.append("")  # Empty line between slides
            
            return "\n".join(text_content), None
        except ImportError:
            return None, "python-pptx not installed. Install with: pip install python-pptx"
        except Exception as exc:
            return None, f"PowerPoint read error: {exc}"